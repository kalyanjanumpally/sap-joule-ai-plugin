"""
Generates the customer pitch deck for @saptarishi/cds-plugin-llm +
@saptarishi/cds-plugin-vector-hana.

Executive overview: 12 slides, 16:9. SAP-friendly navy + coral palette.
Run:  python3 build_deck.py
Out:  ./cds-plugin-llm-and-vector-hana.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---- Palette (SAP-esque: deep navy primary, coral accent) --------------
NAVY = RGBColor(0x0F, 0x25, 0x3C)         # primary background / titles
CORAL = RGBColor(0xE0, 0x4E, 0x39)        # accents / highlights
STEEL = RGBColor(0x4A, 0x6B, 0x8C)        # secondary text
CLOUD = RGBColor(0xF4, 0xF6, 0xF9)        # slide background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1B, 0x1E, 0x24)          # body text on light
MOSS = RGBColor(0x2E, 0x8B, 0x57)         # green for OK / savings numbers
AMBER = RGBColor(0xE0, 0x99, 0x1F)        # amber for warnings / caveats

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_bg(slide, rgb):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb


def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, size=18, color=INK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Calibri'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return tb


def add_bullets(slide, x, y, w, h, items, size=16, color=INK,
                bullet_color=CORAL, spacing=8):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        # bullet marker
        r1 = p.add_run()
        r1.text = "▸ "
        r1.font.name = 'Calibri'
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        # body
        r2 = p.add_run()
        r2.text = item
        r2.font.name = 'Calibri'
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb


def add_footer(slide, page_num, total, title="SAP Joule + CAP-native LLM plugins"):
    add_rect(slide, 0, SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), NAVY)
    add_text(slide, Inches(0.4), SLIDE_H - Inches(0.35), Inches(9),
             Inches(0.35), title, size=10, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.35),
             Inches(0.8), Inches(0.35), f"{page_num} / {total}",
             size=10, color=WHITE, align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.MIDDLE)


def add_slide_title(slide, title, kicker=None):
    # accent bar
    add_rect(slide, Inches(0.5), Inches(0.55), Inches(0.15),
             Inches(0.6), CORAL)
    if kicker:
        add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.35),
                 kicker.upper(), size=11, color=CORAL, bold=True)
        add_text(slide, Inches(0.8), Inches(0.85), Inches(12), Inches(0.6),
                 title, size=28, color=NAVY, bold=True)
    else:
        add_text(slide, Inches(0.8), Inches(0.55), Inches(12), Inches(0.7),
                 title, size=30, color=NAVY, bold=True)


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # Blank
    s = prs.slides.add_slide(layout)
    set_bg(s, CLOUD)
    return s


# ---- Slides -------------------------------------------------------------

def slide_1_title(prs, total):
    s = blank_slide(prs)
    set_bg(s, NAVY)
    # Left coral bar
    add_rect(s, 0, 0, Inches(0.4), SLIDE_H, CORAL)
    # Title
    add_text(s, Inches(1.0), Inches(2.2), Inches(11), Inches(0.5),
             "@SAPTARISHI",  size=14, color=CORAL, bold=True)
    add_text(s, Inches(1.0), Inches(2.6), Inches(11), Inches(1.2),
             "cds-plugin-llm  +  cds-plugin-vector-hana",
             size=38, color=WHITE, bold=True)
    add_text(s, Inches(1.0), Inches(3.9), Inches(11), Inches(0.6),
             "Two CAP-native plugins that turn SAP Joule projects into",
             size=20, color=WHITE)
    add_text(s, Inches(1.0), Inches(4.35), Inches(11), Inches(0.6),
             "production-grade AI platforms — in one line of code.",
             size=20, color=WHITE)
    # Bottom credit line
    add_text(s, Inches(1.0), Inches(6.5), Inches(11), Inches(0.4),
             "cds-plugin-llm 2.38.0  ·  cds-plugin-vector-hana 0.13.0",
             size=12, color=STEEL)
    # No footer on title


def slide_2_problem(prs, page, total):
    s = blank_slide(prs)
    add_slide_title(s, "Every SAP LLM project rebuilds the same 40 layers",
                    kicker="The problem")
    add_bullets(s, Inches(0.8), Inches(1.9), Inches(12), Inches(4.5), [
        "The demo works.  Then finance asks for a budget cap. Ops asks for retries. "
        "Legal asks for PII scrubbing. Security asks for prompt-injection defense.",
        "Every team hand-builds cost / resilience / security / observability from "
        "scratch. Six months later there's still no A/B framework, no rate-limit "
        "retry with jitter, no cross-provider fallback, no cost forecast.",
        "SAP's own Generative AI Hub gives you a model endpoint. It does NOT give "
        "you the 40 middleware layers between the endpoint and a production POS.",
        "Meanwhile CAP is the SAP-native way to build services. What if the "
        "AI plumbing came in the same shape as any other CAP dependency?",
    ], size=16)
    add_footer(s, page, total)


def slide_3_answer(prs, page, total):
    s = blank_slide(prs)
    add_slide_title(s, "Two plugins.  One line each.",
                    kicker="Our answer")

    # Left card: cds-plugin-llm
    left_x = Inches(0.7)
    card_w = Inches(6.0)
    add_rect(s, left_x, Inches(1.8), card_w, Inches(4.9), WHITE)
    add_rect(s, left_x, Inches(1.8), card_w, Inches(0.5), NAVY)
    add_text(s, left_x + Inches(0.2), Inches(1.85), card_w, Inches(0.5),
             "cds-plugin-llm",
             size=18, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left_x + Inches(0.25), Inches(2.5), card_w - Inches(0.5),
             Inches(0.5), "11 providers behind one API",
             size=14, color=STEEL, bold=True)
    add_bullets(s, left_x + Inches(0.25), Inches(3.0),
                card_w - Inches(0.5), Inches(3.5), [
        "Anthropic · Google Gemini · Bedrock · Azure OpenAI",
        "Ollama · Groq · Fireworks · DeepSeek · Mistral",
        "OpenAI-compatible · SAP Generative AI Hub",
        "40+ middleware primitives · 2.x API stability contract",
        "Koa-style llm.use() composition",
    ], size=14)

    # Right card: vector-hana
    right_x = Inches(6.9)
    add_rect(s, right_x, Inches(1.8), card_w, Inches(4.9), WHITE)
    add_rect(s, right_x, Inches(1.8), card_w, Inches(0.5), NAVY)
    add_text(s, right_x + Inches(0.2), Inches(1.85), card_w, Inches(0.5),
             "cds-plugin-vector-hana",
             size=18, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, right_x + Inches(0.25), Inches(2.5), card_w - Inches(0.5),
             Inches(0.5), "@rag annotation on any CAP entity",
             size=14, color=STEEL, bold=True)
    add_bullets(s, right_x + Inches(0.25), Inches(3.0),
                card_w - Inches(0.5), Inches(3.5), [
        "HANA Cloud REAL_VECTOR + SQLite dev fallback",
        "Hybrid vector + keyword (RRF fusion)",
        "HyDE query expansion, LLM rerank",
        "Auto-declares searchByMeaning / askAbout OData actions",
        "Zero handler code — pure declarative RAG",
    ], size=14)

    add_footer(s, page, total)


def slide_4_llm_categories(prs, page, total):
    s = blank_slide(prs)
    add_slide_title(s, "cds-plugin-llm — 40+ primitives across 10 categories",
                    kicker="Plugin 1")

    cats = [
        ("Cost", CORAL, "budget · guard · forecast · quota · aware routing · overrun predictor"),
        ("Resilience", RGBColor(0x2E, 0x8B, 0x57), "circuit breaker · bulkhead · deadline · retry · region failover · hedge"),
        ("Security", RGBColor(0x8E, 0x44, 0xAD), "PII redact · reversible tokens · injection guard · guardrails · HMAC signing"),
        ("Observability", RGBColor(0x0F, 0x6E, 0x99), "OTel · Prometheus · JSON log · health probe · latency histograms · MCP"),
        ("Routing", RGBColor(0xE0, 0x99, 0x1F), "model · tenant · cost-aware · fair-share · semantic · load balancer"),
        ("Caching", RGBColor(0xC0, 0x39, 0x2B), "exact · semantic · embedding dedup · request coalescer · fuzzy dedup"),
        ("Contract", RGBColor(0x16, 0xA0, 0x85), "JSON schema validate · repair · function-call arbitrator · length gate"),
        ("Long-context", RGBColor(0x5D, 0x6D, 0x7E), "compact history · session store · reversible tokenization"),
        ("RAG + Eval", RGBColor(0xD3, 0x50, 0x0C), "ragChain · llmJudge · promptRegression · consensusVoting · A/B experiment"),
        ("Multi-agent", RGBColor(0x7D, 0x3C, 0x98), "Agent · runAgents · streamAgents · autoToolChain · supervisor"),
    ]

    # 5 x 2 grid
    top = Inches(1.85)
    row_h = Inches(1.05)
    col_w = Inches(6.15)
    gap_x = Inches(0.15)
    left_x = Inches(0.5)
    for i, (name, color, desc) in enumerate(cats):
        col = i % 2
        row = i // 2
        x = left_x + (col_w + gap_x) * col
        y = top + row_h * row
        add_rect(s, x, y, col_w, Inches(0.95), WHITE)
        add_rect(s, x, y, Inches(0.15), Inches(0.95), color)
        add_text(s, x + Inches(0.25), y + Inches(0.05), Inches(2.2),
                 Inches(0.4), name, size=15, color=NAVY, bold=True)
        add_text(s, x + Inches(0.25), y + Inches(0.42), col_w - Inches(0.35),
                 Inches(0.5), desc, size=11, color=INK)

    add_footer(s, page, total)


def slide_5_headline_costs(prs, page, total):
    s = blank_slide(prs)
    add_slide_title(s, "Cost primitives — turn AI spend into an SLO",
                    kicker="Business case")

    # Split: 2/3 left bullet list, 1/3 right callout
    add_bullets(s, Inches(0.8), Inches(1.9), Inches(7.5), Inches(4.5), [
        "costGuard — pre-flight per-call USD ceiling (rejects before the API call)",
        "costBudget — per-tenant / per-model / per-window ceilings; hydrate from a CAP entity, edit live, no restart",
        "costForecast — rolling-window spend projection with rising-edge alerts (\"you'll hit the $50 target at 2:47pm\")",
        "quotaManager — per-user USD quota with warnings before the ceiling",
        "costOverrunPredictor — calendar-window burn-rate + startOfMonth / endOfMonth helpers",
        "costAwareRouter — cheap model first, escalate to expensive on schema failure",
        "adaptiveMaxTokens — shrink maxTokens automatically to fit remaining budget",
    ], size=14)

    # Callout
    card_x = Inches(8.6)
    card_w = Inches(4.3)
    add_rect(s, card_x, Inches(2.0), card_w, Inches(4.5), NAVY)
    add_text(s, card_x + Inches(0.25), Inches(2.15), card_w - Inches(0.5),
             Inches(0.5), "WHY IT MATTERS", size=11, color=CORAL, bold=True)
    add_text(s, card_x + Inches(0.25), Inches(2.55), card_w - Inches(0.5),
             Inches(1.0), "The finance conversation flips.",
             size=20, color=WHITE, bold=True)
    add_text(s, card_x + Inches(0.25), Inches(3.4), card_w - Inches(0.5),
             Inches(2.5),
             "Instead of \"what did this cost last month?\" (reactive), "
             "finance now says \"here's the ceiling per tenant per hour\" "
             "and the plumbing enforces it. Overspend becomes a "
             "structurally impossible failure mode, not a monitoring problem.",
             size=13, color=WHITE)

    add_footer(s, page, total)


def slide_6_resilience_security(prs, page, total):
    s = blank_slide(prs)
    add_slide_title(s, "Resilience + Security — production defaults, not homework",
                    kicker="Ops + Security case")

    # Two columns
    col1_x = Inches(0.7)
    col2_x = Inches(7.0)
    col_w = Inches(5.9)

    # Resilience header
    add_rect(s, col1_x, Inches(1.85), col_w, Inches(0.5), NAVY)
    add_text(s, col1_x + Inches(0.2), Inches(1.9), col_w, Inches(0.5),
             "Resilience", size=16, color=WHITE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, col1_x + Inches(0.15), Inches(2.5), col_w - Inches(0.3),
                Inches(4.4), [
        "retryOnRateLimit — exponential backoff, honors Retry-After",
        "circuitBreaker — per-provider open/half-open/closed",
        "bulkhead + adaptiveBulkhead — AIMD tuning on p95 latency",
        "deadline — hard per-call time budget",
        "gracePeriod — soft deadline warnings + optional hard timeout",
        "speculativeHedge — staggered parallel hedges for tail latency",
        "regionFailover — chatWithFallback across N regions",
        "chaosInjector — deterministic seeded fault injection (test-only)",
    ], size=13, spacing=5)

    # Security header
    add_rect(s, col2_x, Inches(1.85), col_w, Inches(0.5), NAVY)
    add_text(s, col2_x + Inches(0.2), Inches(1.9), col_w, Inches(0.5),
             "Security", size=16, color=WHITE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, col2_x + Inches(0.15), Inches(2.5), col_w - Inches(0.3),
                Inches(4.4), [
        "guardrails — regex + PII + blocklist stages",
        "promptInjectionGuard — 6-detector jailbreak defense",
        "piiRedact — round-trip PII masking with reversible tokens",
        "safetyClassifier — moderation + Anthropic refusal detection",
        "requestSigning — HMAC receipts + verifyReceiptChain",
        "responseSigning — HMAC responses (tamper-evident audit trail)",
        "sensitiveDataAudit — before/after diff capture",
        "emptyResponseDetector — refusal-pattern detection + auto-retry",
    ], size=13, spacing=5)

    add_footer(s, page, total)


def slide_7_observability(prs, page, total):
    s = blank_slide(prs)
    add_slide_title(s, "Observability — every middleware is a live MCP resource",
                    kicker="SRE case")

    add_bullets(s, Inches(0.8), Inches(1.9), Inches(7.5), Inches(4.5), [
        "OpenTelemetry spans with cost + correlation + routing + errors",
        "Prometheus /metrics text-exposition — one endpoint, all counters",
        "JSON structured logs (one line per call) with error taxonomy",
        "providerHealthProbe + healthAggregate — unified verdict per provider",
        "latencyHistogram — per-dimension p50/p95/p99 with Prom export",
        "traceCorrelation via uuidv7 (sortable, deterministic)",
        "replayBuffer — rolling in-memory window for on-call debugging",
        "usageMeteringToCap — auto-persists to a CAP entity (queryable via OData)",
    ], size=14)

    # Right callout with MCP flavour
    card_x = Inches(8.6)
    card_w = Inches(4.3)
    add_rect(s, card_x, Inches(2.0), card_w, Inches(4.5), NAVY)
    add_text(s, card_x + Inches(0.25), Inches(2.15), card_w - Inches(0.5),
             Inches(0.5), "MCP-FIRST", size=11, color=CORAL, bold=True)
    add_text(s, card_x + Inches(0.25), Inches(2.55), card_w - Inches(0.5),
             Inches(1.0), "Every counter is a subscribable resource.",
             size=18, color=WHITE, bold=True)
    add_text(s, card_x + Inches(0.25), Inches(3.7), card_w - Inches(0.5),
             Inches(2.5),
             "asMcpResource() ships on every middleware. Claude Desktop, "
             "Cline, Cursor can subscribe to config://budget, "
             "config://circuit-breaker, config://fuzzy-dedup — and get a "
             "push notification the moment a counter changes.",
             size=13, color=WHITE)

    add_footer(s, page, total)


def slide_8_fuzzy_dedup(prs, page, total):
    s = blank_slide(prs)
    add_slide_title(s, "Fuzzy Dedup — the third layer in the dedup story",
                    kicker="New in 2.38.0")

    # Left: bullets
    add_bullets(s, Inches(0.7), Inches(1.9), Inches(6.5), Inches(3.5), [
        "Layer 1: requestCoalescer — byte-identical, in-flight",
        "Layer 2: responseCache — byte-identical, at-rest",
        "Layer 3: semanticCache — embedding-similar (needs embedder)",
        "Layer 4 (NEW): fuzzyDedup — character-similar (no embedder)",
    ], size=15)

    add_text(s, Inches(0.7), Inches(4.3), Inches(6.5), Inches(1.0),
             "Jaccard-trigram or normalized Levenshtein similarity. "
             "Catches supplier IDs retyped with dashes vs spaces, "
             "whitespace-only diffs, single-character typos, or extra "
             "optional fields — BEFORE the semantic cache spends an "
             "embedding round-trip.",
             size=12, color=INK)

    # Right: proof point card
    card_x = Inches(7.7)
    card_w = Inches(5.2)
    add_rect(s, card_x, Inches(1.9), card_w, Inches(4.8), WHITE,
             line=STEEL)
    add_rect(s, card_x, Inches(1.9), card_w, Inches(0.55), CORAL)
    add_text(s, card_x + Inches(0.2), Inches(1.95), card_w - Inches(0.4),
             Inches(0.55), "LIVE ON THE PROCUREMENT COPILOT",
             size=12, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, card_x + Inches(0.3), Inches(2.7), card_w - Inches(0.6),
             Inches(0.4), "3 near-duplicate POST /ai/summarizePurchaseOrder calls",
             size=12, color=STEEL, bold=True)

    stats = [
        ("totalCalls", "3", INK),
        ("fuzzyHits", "2", MOSS),
        ("hitRate", "0.667", MOSS),
        ("lastSimilarity", "0.834", INK),
        ("LLM calls saved", "2 of 3", MOSS),
    ]
    y_top = Inches(3.2)
    for i, (k, v, c) in enumerate(stats):
        y = y_top + Inches(0.55) * i
        add_text(s, card_x + Inches(0.3), y, Inches(2.5), Inches(0.4),
                 k, size=13, color=STEEL)
        add_text(s, card_x + Inches(2.8), y, Inches(2.3), Inches(0.4),
                 v, size=15, color=c, bold=True)

    add_footer(s, page, total)


def slide_9_vector_hana(prs, page, total):
    s = blank_slide(prs)
    add_slide_title(s, "cds-plugin-vector-hana — @rag on any CAP entity",
                    kicker="Plugin 2")

    # Left: bullets
    add_bullets(s, Inches(0.7), Inches(1.9), Inches(6.5), Inches(4.5), [
        "One annotation on your existing CAP entity — no boilerplate handlers",
        "HANA Cloud REAL_VECTOR + COSINE_SIMILARITY in production",
        "SQLite fallback for local dev — no HANA needed to build features",
        "Hybrid search (vector + keyword via FTS5 / CONTAINS) with alpha knob",
        "Metadata filters: filter: { region: 'EMEA' } exact-match on any field",
        "Auto-declares searchByMeaning + askAbout OData actions on the entity",
        "Composes with cds-plugin-llm for embeddings + rerank + answer",
    ], size=14)

    # Right: code card
    card_x = Inches(7.7)
    card_w = Inches(5.2)
    add_rect(s, card_x, Inches(1.9), card_w, Inches(4.8), NAVY)
    add_text(s, card_x + Inches(0.25), Inches(2.0), card_w - Inches(0.5),
             Inches(0.4), "CDS FILE", size=11, color=CORAL, bold=True)

    code = ("entity SupplierContracts {\n"
            "  key ID   : String;\n"
            "  supplier : String;\n"
            "  region   : String;\n"
            "  terms    : LargeString;\n"
            "} @rag: {\n"
            "  fields:         ['supplier','region','terms'],\n"
            "  metadataFields: ['region','supplier'],\n"
            "  dimension:      768,\n"
            "  store:          'hana',\n"
            "  search:         'hybrid',\n"
            "};")

    tb = s.shapes.add_textbox(card_x + Inches(0.25), Inches(2.4),
                              card_w - Inches(0.5), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = 'Consolas'
        r.font.size = Pt(12)
        r.font.color.rgb = WHITE

    add_footer(s, page, total)


def slide_10_rag_pipeline(prs, page, total):
    s = blank_slide(prs)
    add_slide_title(s, "5-stage RAG pipeline — every stage a separate primitive",
                    kicker="cds-plugin-vector-hana in action")

    stages = [
        ("1. HyDE expand",
         "LLM writes a hypothetical answer — both question + answer get embedded"),
        ("2. Hybrid retrieve",
         "Vector + keyword search across each expanded query in parallel"),
        ("3. RRF fusion",
         "Reciprocal Rank Fusion — docs in multiple lists rank higher"),
        ("4. LLM rerank",
         "Structured 0-10 relevance scoring rewrites the top-K order"),
        ("5. RAG answer",
         "Augmented prompt with cited sources → chat completion"),
    ]

    # 5 horizontal cards
    total_w = SLIDE_W - Inches(1.4)
    card_w = (total_w - Inches(0.4)) / 5
    left = Inches(0.7)
    top = Inches(2.2)
    card_h = Inches(3.0)
    for i, (title_, desc) in enumerate(stages):
        x = left + (card_w + Inches(0.1)) * i
        add_rect(s, x, top, card_w, card_h, WHITE, line=STEEL)
        add_rect(s, x, top, card_w, Inches(0.6), NAVY)
        add_text(s, x, top, card_w, Inches(0.6), title_,
                 size=13, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.1), top + Inches(0.75),
                 card_w - Inches(0.2), card_h - Inches(0.8), desc,
                 size=11, color=INK)

    add_text(s, Inches(0.7), Inches(5.7), Inches(12), Inches(1.0),
             "Groq Llama-3.3-70B end-to-end: ~2.5s.  Every stage is "
             "individually swappable (custom rerankers, alternate expanders, "
             "different fusion strategies) — this is the composed default.",
             size=13, color=STEEL, align=PP_ALIGN.CENTER)

    add_footer(s, page, total)


def slide_11_proof_point(prs, page, total):
    s = blank_slide(prs)
    add_slide_title(s, "Joule Procurement Copilot — real customer-ready proof",
                    kicker="Proof point")

    # Left: what
    add_bullets(s, Inches(0.7), Inches(1.9), Inches(7.5), Inches(4.5), [
        "A live SAP Joule agent using both plugins end-to-end",
        "Actions: summarize PO · explain invoice risk · extract line items · "
        "multi-agent scenario analysis · voice memo → PO draft · batch supplier scoring",
        "29-layer middleware chain (OTel → fuzzy dedup → semantic cache → provider)",
        "Auto-declared RAG actions on SupplierContracts entity (hybrid + HyDE)",
        "Deployed live on Render for the demo call — public URL, click and demo",
        "MCP observability surface on port 3334 — subscribe to any middleware's counters",
    ], size=13, spacing=6)

    # Right: numbers card
    card_x = Inches(8.6)
    card_w = Inches(4.3)
    add_rect(s, card_x, Inches(1.9), card_w, Inches(4.8), NAVY)

    numbers = [
        ("29", "middleware layers", WHITE),
        ("11", "LLM providers", WHITE),
        ("40+", "primitives shipped", WHITE),
        ("~2.5s", "5-stage RAG p95", CORAL),
        ("0.667", "fuzzy hit rate (demo)", MOSS),
    ]
    y = Inches(2.15)
    for num, lbl, c in numbers:
        add_text(s, card_x + Inches(0.3), y, card_w - Inches(0.6),
                 Inches(0.7), num, size=30, color=c, bold=True)
        add_text(s, card_x + Inches(0.3), y + Inches(0.6),
                 card_w - Inches(0.6), Inches(0.35), lbl,
                 size=12, color=WHITE)
        y = y + Inches(0.95)

    add_footer(s, page, total)


def slide_12_getting_started(prs, page, total):
    s = blank_slide(prs)
    add_slide_title(s, "Getting started — one npm install, one line of code",
                    kicker="Next")

    # Code card
    code_x = Inches(0.7)
    code_w = Inches(7.5)
    add_rect(s, code_x, Inches(1.9), code_w, Inches(4.5), NAVY)
    add_text(s, code_x + Inches(0.25), Inches(2.0), code_w - Inches(0.5),
             Inches(0.4), "TERMINAL", size=11, color=CORAL, bold=True)

    code = ("$ npm install @saptarishi/cds-plugin-llm \\\n"
            "              @saptarishi/cds-plugin-vector-hana\n"
            "\n"
            "// srv/ai-service.js\n"
            "const llm = await cds.connect.to('llm');\n"
            "\n"
            "llm.use(costBudget({ perTenant: { acme: 100 } }));\n"
            "llm.use(promptInjectionGuard());\n"
            "llm.use(piiRedact());\n"
            "llm.use(fuzzyDedup({ store: inMemoryFuzzyStore() }));\n"
            "llm.use(responseCache({ semantic: { embedder } }));\n"
            "\n"
            "const { text } = await llm.chat({\n"
            "  messages: [{ role: 'user', content: 'summarise PO-4471' }],\n"
            "});")

    tb = s.shapes.add_textbox(code_x + Inches(0.25), Inches(2.4),
                              code_w - Inches(0.5), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = 'Consolas'
        r.font.size = Pt(11)
        r.font.color.rgb = WHITE

    # Right: links
    card_x = Inches(8.6)
    card_w = Inches(4.3)
    add_rect(s, card_x, Inches(1.9), card_w, Inches(4.5), WHITE,
             line=STEEL)
    add_text(s, card_x + Inches(0.25), Inches(2.05), card_w - Inches(0.5),
             Inches(0.5), "RESOURCES", size=12, color=CORAL, bold=True)

    add_text(s, card_x + Inches(0.25), Inches(2.6), card_w - Inches(0.5),
             Inches(0.35), "npm", size=13, color=STEEL, bold=True)
    add_text(s, card_x + Inches(0.25), Inches(2.9), card_w - Inches(0.5),
             Inches(0.35),
             "@saptarishi/cds-plugin-llm@2.38.0",
             size=11, color=INK)
    add_text(s, card_x + Inches(0.25), Inches(3.2), card_w - Inches(0.5),
             Inches(0.35),
             "@saptarishi/cds-plugin-vector-hana@0.13.0",
             size=11, color=INK)

    add_text(s, card_x + Inches(0.25), Inches(3.75), card_w - Inches(0.5),
             Inches(0.35), "Source", size=13, color=STEEL, bold=True)
    add_text(s, card_x + Inches(0.25), Inches(4.05), card_w - Inches(0.5),
             Inches(0.5),
             "github.com/kalyanjanumpally/\nsap-joule-ai-plugin",
             size=11, color=INK)

    add_text(s, card_x + Inches(0.25), Inches(4.85), card_w - Inches(0.5),
             Inches(0.35), "Live demo", size=13, color=STEEL, bold=True)
    add_text(s, card_x + Inches(0.25), Inches(5.15), card_w - Inches(0.5),
             Inches(0.7),
             "Procurement Copilot on Render (URL provided during the call)",
             size=11, color=INK)

    add_text(s, card_x + Inches(0.25), Inches(5.85), card_w - Inches(0.5),
             Inches(0.35), "Contact", size=13, color=STEEL, bold=True)
    add_text(s, card_x + Inches(0.25), Inches(6.15), card_w - Inches(0.5),
             Inches(0.35), "jkalyan@alumni.iitm.ac.in", size=11, color=INK)

    add_footer(s, page, total)


# ---- Build --------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = [
        slide_1_title,
        slide_2_problem,
        slide_3_answer,
        slide_4_llm_categories,
        slide_5_headline_costs,
        slide_6_resilience_security,
        slide_7_observability,
        slide_8_fuzzy_dedup,
        slide_9_vector_hana,
        slide_10_rag_pipeline,
        slide_11_proof_point,
        slide_12_getting_started,
    ]
    total = len(slides)
    for i, fn in enumerate(slides, start=1):
        if fn is slide_1_title:
            fn(prs, total)
        else:
            fn(prs, i, total)

    out = "cds-plugin-llm-and-vector-hana.pptx"
    prs.save(out)
    print(f"Wrote {out}  ({total} slides)")


if __name__ == "__main__":
    build()
